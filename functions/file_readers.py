# -*- coding: utf-8 -*-
"""
Leitores especializados para diferentes tipos de arquivo
Cada leitor extrai conteúdo relevante para gerar nomes de arquivo
"""

import os
from pathlib import Path
from typing import Optional, Dict, Any
import re

# Importações condicionais para evitar erros se alguma biblioteca não estiver instalada
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import PyPDF2
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

from .encoding_utils import normalize_text, safe_filename

class FileReader:
    """Classe base para leitores de arquivo"""
    
    def __init__(self):
        self.max_title_length = 50
        self.max_content_chars = 1000
    
    def read_file(self, file_path: str) -> Dict[str, Any]:
        """
        Lê um arquivo e extrai informações relevantes
        
        Returns:
            {
                'title': str,           # Título/nome sugerido
                'content_preview': str, # Preview do conteúdo
                'success': bool,        # Se a leitura foi bem-sucedida
                'error': str,          # Mensagem de erro se houver
                'suggested_name': str   # Nome de arquivo sugerido
            }
        """
        try:
            result = self._extract_content(file_path)
            if result['success']:
                result['suggested_name'] = self._generate_filename(result['title'], file_path)
            return result
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler arquivo: {str(e)}",
                'suggested_name': ''
            }
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        """Método a ser implementado pelas subclasses"""
        raise NotImplementedError
    
    def _generate_filename(self, title: str, original_path: str) -> str:
        """Gera nome de arquivo baseado no título extraído - apenas para compatibilidade"""
        # Este método não é mais usado, pois a geração de nomes foi movida para FileRenamer
        return Path(original_path).name

class WordReader(FileReader):
    """Leitor para arquivos Word (.docx)"""
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        if not DOCX_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Biblioteca python-docx não disponível'
            }
        
        try:
            doc = Document(file_path)
            
            # Tentar extrair título de várias formas
            title = ""
            content_parts = []
            
            # 1. Verificar propriedades do documento
            if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
                title = doc.core_properties.title
            
            # 2. Se não tem título, usar primeiro parágrafo
            if not title and doc.paragraphs:
                first_paragraph = doc.paragraphs[0].text.strip()
                if first_paragraph:
                    title = first_paragraph
            
            # 3. Coletar conteúdo para preview
            for paragraph in doc.paragraphs[:5]:  # Primeiros 5 parágrafos
                text = paragraph.text.strip()
                if text:
                    content_parts.append(text)
            
            content_preview = ' '.join(content_parts)[:self.max_content_chars]
            
            return {
                'title': title or "documento_word",
                'content_preview': content_preview,
                'success': True,
                'error': ''
            }
            
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler Word: {str(e)}"
            }

class ExcelReader(FileReader):
    """Leitor para arquivos Excel (.xlsx e .xls)"""
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        file_extension = Path(file_path).suffix.lower()
        
        # Tentar .xlsx primeiro com openpyxl
        if file_extension == '.xlsx':
            return self._read_xlsx(file_path)
        
        # Para .xls usar xlrd
        elif file_extension == '.xls':
            return self._read_xls(file_path)
        
        else:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f'Extensão não suportada: {file_extension}'
            }
    
    def _read_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Lê arquivos .xlsx com openpyxl"""
        if not OPENPYXL_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Biblioteca openpyxl não disponível'
            }
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            # Usar nome da primeira planilha como título
            sheet_names = workbook.sheetnames
            title = sheet_names[0] if sheet_names else "planilha"
            
            # Se o nome é genérico, tentar usar conteúdo da célula A1
            if title.lower() in ['sheet1', 'planilha1', 'plan1']:
                try:
                    first_sheet = workbook.active
                    cell_a1 = first_sheet['A1'].value
                    if cell_a1 and isinstance(cell_a1, str):
                        title = str(cell_a1).strip()
                except:
                    pass
            
            # Coletar preview do conteúdo
            content_parts = []
            try:
                sheet = workbook.active
                for row in range(1, min(6, sheet.max_row + 1)):  # Primeiras 5 linhas
                    row_data = []
                    for col in range(1, min(4, sheet.max_column + 1)):  # Primeiras 3 colunas
                        cell_value = sheet.cell(row=row, column=col).value
                        if cell_value is not None:
                            row_data.append(str(cell_value))
                    if row_data:
                        content_parts.append(' | '.join(row_data))
            except:
                pass
            
            content_preview = '\n'.join(content_parts)[:self.max_content_chars]
            
            return {
                'title': title,
                'content_preview': content_preview,
                'success': True,
                'error': ''
            }
            
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler .xlsx: {str(e)}"
            }
    
    def _read_xls(self, file_path: str) -> Dict[str, Any]:
        """Lê arquivos .xls legados com xlrd"""
        if not XLRD_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Biblioteca xlrd não disponível para arquivos .xls'
            }
        
        try:
            # Abrir workbook com xlrd
            workbook = xlrd.open_workbook(file_path)
            
            # Usar nome da primeira planilha como título
            sheet_names = workbook.sheet_names()
            title = sheet_names[0] if sheet_names else "planilha"
            
            # Se o nome é genérico, tentar usar conteúdo da célula A1
            if title.lower() in ['sheet1', 'planilha1', 'plan1', 'sheet']:
                try:
                    first_sheet = workbook.sheet_by_index(0)
                    if first_sheet.nrows > 0 and first_sheet.ncols > 0:
                        cell_a1 = first_sheet.cell_value(0, 0)  # row 0, col 0 = A1
                        if cell_a1 and isinstance(cell_a1, str) and cell_a1.strip():
                            title = str(cell_a1).strip()
                except:
                    pass
            
            # Coletar preview do conteúdo
            content_parts = []
            try:
                sheet = workbook.sheet_by_index(0)
                max_rows = min(5, sheet.nrows)
                max_cols = min(3, sheet.ncols)
                
                for row in range(max_rows):
                    row_data = []
                    for col in range(max_cols):
                        try:
                            cell_value = sheet.cell_value(row, col)
                            if cell_value is not None and str(cell_value).strip():
                                # Converter diferentes tipos de célula
                                if isinstance(cell_value, float) and cell_value.is_integer():
                                    row_data.append(str(int(cell_value)))
                                else:
                                    row_data.append(str(cell_value))
                        except:
                            continue
                    
                    if row_data:
                        content_parts.append(' | '.join(row_data))
            except:
                pass
            
            content_preview = '\n'.join(content_parts)[:self.max_content_chars]
            
            return {
                'title': title,
                'content_preview': content_preview,
                'success': True,
                'error': ''
            }
            
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler .xls: {str(e)}"
            }

class PowerPointReader(FileReader):
    """Leitor para arquivos PowerPoint (.pptx)"""
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        if not PPTX_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Biblioteca python-pptx não disponível'
            }
        
        try:
            presentation = Presentation(file_path)
            
            title = ""
            content_parts = []
            
            # Tentar extrair título do primeiro slide
            if presentation.slides:
                first_slide = presentation.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not title:
                            title = shape.text.strip()
                        content_parts.append(shape.text.strip())
            
            # Se não encontrou título, usar "apresentacao"
            if not title:
                title = "apresentacao"
            
            content_preview = ' '.join(content_parts)[:self.max_content_chars]
            
            return {
                'title': title,
                'content_preview': content_preview,
                'success': True,
                'error': ''
            }
            
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler PowerPoint: {str(e)}"
            }

class PDFReader(FileReader):
    """Leitor para arquivos PDF"""
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        if not PDF_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Bibliotecas PDF não disponíveis'
            }
        
        title = ""
        content_preview = ""
        
        # Tentar primeiro com pdfplumber (melhor para texto)
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                if pdf.pages:
                    first_page = pdf.pages[0]
                    text = first_page.extract_text()
                    if text:
                        lines = text.split('\n')
                        # Usar primeira linha não vazia como título
                        for line in lines:
                            line = line.strip()
                            if line and len(line) > 3:
                                title = line
                                break
                        content_preview = text[:self.max_content_chars]
        except:
            pass
        
        # Se pdfplumber falhou, tentar PyPDF2
        if not title:
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    if reader.pages:
                        page = reader.pages[0]
                        text = page.extract_text()
                        if text:
                            lines = text.split('\n')
                            for line in lines:
                                line = line.strip()
                                if line and len(line) > 3:
                                    title = line
                                    break
                            content_preview = text[:self.max_content_chars]
            except:
                pass
        
        return {
            'title': title or "documento_pdf",
            'content_preview': content_preview,
            'success': bool(title or content_preview),
            'error': '' if (title or content_preview) else 'Não foi possível extrair texto do PDF'
        }

class CSVReader(FileReader):
    """Leitor para arquivos CSV"""
    
    def _extract_content(self, file_path: str) -> Dict[str, Any]:
        if not PANDAS_AVAILABLE:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': 'Biblioteca pandas não disponível'
            }
        
        try:
            # Tentar diferentes codificações
            encodings = ['utf-8', 'latin-1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, nrows=5)
                    break
                except:
                    continue
            
            if df is None:
                return {
                    'title': '',
                    'content_preview': '',
                    'success': False,
                    'error': 'Não foi possível ler o arquivo CSV'
                }
            
            # Usar nome das colunas como base do título
            columns = list(df.columns)
            if columns:
                title = '_'.join(columns[:3])  # Primeiras 3 colunas
            else:
                title = "dados_csv"
            
            # Preview do conteúdo
            content_preview = df.head().to_string()[:self.max_content_chars]
            
            return {
                'title': title,
                'content_preview': content_preview,
                'success': True,
                'error': ''
            }
            
        except Exception as e:
            return {
                'title': '',
                'content_preview': '',
                'success': False,
                'error': f"Erro ao ler CSV: {str(e)}"
            }

# Factory para criar leitores
def get_file_reader(file_type: str) -> Optional[FileReader]:
    """Retorna o leitor apropriado para o tipo de arquivo"""
    readers = {
        'word': WordReader,
        'excel': ExcelReader,
        'powerpoint': PowerPointReader,
        'pdf': PDFReader,
        'csv': CSVReader
    }
    
    reader_class = readers.get(file_type)
    return reader_class() if reader_class else None